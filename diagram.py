# -*- coding: utf-8 -*-
"""금융구조도 — 틀(PPT)에서 입력칸을 뽑고, 채운 구조도를 만든다.

레이아웃 틀: layout/금융구조도_모음.pptx  (한 슬라이드 = 구조도 하나)

틀에 붙은 번호(①, ①-1, ②…)는 '이 칸이 뭐냐'를 알려주려고 달아둔 표식이다.
  · 상자(표) 옆에 따로 떠 있는 동그라미 숫자  → 그 상자 칸의 이름
  · 화살표 라벨은 글 앞머리에 숫자가 붙어 있음 → 그 라벨의 이름
회사 원본 구조도에는 번호가 없으므로 **완성된 구조도에서는 지운다.**
"""
import os
import re
from copy import deepcopy

from pptx import Presentation
from pptx.util import Emu, Pt

_HERE = os.path.dirname(os.path.abspath(__file__))
LAYOUT = os.path.join(_HERE, "layout", "금융구조도_모음.pptx")
THUMBS = os.path.join(_HERE, "layout", "구조도_썸네일")

# 동그라미 숫자 ①~⑳ (그 뒤에 -1 같은 꼬리표가 붙을 수 있다)
_CIRCLE = "①②③④⑤⑥⑦⑧⑨⑩⑪⑫⑬⑭⑮⑯⑰⑱⑲⑳"
_MARK_ONLY = re.compile(r"^\s*([" + _CIRCLE + r"])\s*(-\s*\d+)?\s*$")
_MARK_HEAD = re.compile(r"^\s*([" + _CIRCLE + r"])\s*(-\s*\d+)?\s*")

EMU_IN = 914400.0


def _norm_mark(ch, tail):
    """'①' + '-1' → '①-1' (공백 제거)"""
    return ch + (tail.replace(" ", "") if tail else "")


# ── 도형 훑기 ────────────────────────────────────────
def _walk(shapes, path=()):
    """(도형, 인덱스경로) 를 준다. 그룹 안쪽까지 들어간다."""
    for i, sh in enumerate(shapes):
        p = path + (i,)
        yield sh, p
        if sh.shape_type == 6:          # GROUP
            yield from _walk(sh.shapes, p)


def _by_path(slide, path):
    """인덱스경로로 도형을 되찾는다."""
    cur = slide.shapes
    sh = None
    for i in path:
        sh = cur[i]
        if sh.shape_type == 6:
            cur = sh.shapes
    return sh


def _abs_box(sh):
    """도형의 화면상 위치(인치). 그룹 안이면 그룹 변환을 반영한다."""
    try:
        l, t = (sh.left or 0) / EMU_IN, (sh.top or 0) / EMU_IN
        w, h = (sh.width or 0) / EMU_IN, (sh.height or 0) / EMU_IN
        return l, t, w, h
    except Exception:
        return 0.0, 0.0, 0.0, 0.0


def _cell_box(shape, table, r, c):
    """표 한 칸의 위치(인치) — 행높이·열너비를 더해서 구한다."""
    l = (shape.left or 0) / EMU_IN
    t = (shape.top or 0) / EMU_IN
    for j in range(c):
        l += (table.columns[j].width or 0) / EMU_IN
    for i in range(r):
        t += (table.rows[i].height or 0) / EMU_IN
    w = (table.columns[c].width or 0) / EMU_IN
    h = (table.rows[r].height or 0) / EMU_IN
    return l, t, w, h


# ── 구조도 읽기 ──────────────────────────────────────
def read_layouts(path: str = None):
    """틀 PPT 를 읽어 구조도 목록을 만든다.

    반환: [{ "no": 1, "title": "1. 담보대출 Tr. B",
             "fields": [ {...}, ... ] }, ...]

    field = {
      "key":   "box:①-1" 처럼 고유 이름
      "kind":  "상자" | "화살표"
      "mark":  "①-1"  (없으면 "")
      "text":  틀에 적힌 지금 글자
      "loc":   ("cell", path, r, c) | ("tb", path)
    }
    """
    prs = Presentation(path or LAYOUT)
    out = []
    for si, slide in enumerate(prs.slides, 1):
        marks, cells, boxes, title = [], [], [], ""
        for sh, p in _walk(slide.shapes):
            if sh.has_table:
                t = sh.table
                for r in range(len(t.rows)):
                    for c in range(len(t.columns)):
                        txt = (t.cell(r, c).text or "").strip()
                        if txt:
                            cells.append((_cell_box(sh, t, r, c), txt, p, r, c))
                continue
            if not sh.has_text_frame:
                continue
            txt = (sh.text_frame.text or "").strip()
            if not txt:
                continue
            m = _MARK_ONLY.match(txt)
            if m:                                   # 표식만 있는 작은 글상자
                marks.append((_abs_box(sh), _norm_mark(m.group(1), m.group(2))))
                continue
            if not title and _abs_box(sh)[1] < 0.9:  # 맨 위 = 구조도 제목
                title = txt
                continue
            boxes.append((_abs_box(sh), txt, p))

        # 이름 붙일 대상 모으기 — 글 앞머리에 번호가 있으면 그걸 그대로 쓴다.
        fields, need = [], []
        for box, txt, p, r, c in sorted(cells, key=lambda x: (x[0][1], x[0][0])):
            f = {"kind": "상자", "mark": "", "text": txt, "loc": ("cell", p, r, c)}
            m = _MARK_HEAD.match(txt)
            if m and len(txt) > m.end():
                f["mark"] = _norm_mark(m.group(1), m.group(2))
                f["text"] = txt[m.end():].strip()
            else:
                need.append((len(fields), box))
            fields.append(f)
        for box, txt, p in sorted(boxes, key=lambda x: (x[0][1], x[0][0])):
            f = {"kind": "화살표", "mark": "", "text": txt, "loc": ("tb", p)}
            m = _MARK_HEAD.match(txt)
            if m and len(txt) > m.end():
                f["mark"] = _norm_mark(m.group(1), m.group(2))
                f["text"] = txt[m.end():].strip()
            else:
                need.append((len(fields), box))
            fields.append(f)

        # ★남은 칸에 표식을 붙인다. 칸마다 차례로 '제일 가까운 것'을 집으면
        #   먼저 나온 칸이 남의 표식을 가져가 뒤가 밀린다(1·2번 구조도에서 실제로 그랬다).
        #   → 모든 (표식, 칸) 짝의 거리를 다 재서 **가까운 짝부터** 확정한다.
        #   거리는 중심끼리가 아니라 '상자 테두리까지'로 잰다. 표식은 상자 바로
        #   바깥에 붙어 있어서 테두리 거리가 훨씬 잘 갈린다.
        def edge_dist(mb, box):
            ml, mt, mw, mh = mb
            mx, my = ml + mw / 2, mt + mh / 2
            l, t, w, h = box
            dx = max(l - mx, 0, mx - (l + w))
            dy = max(t - my, 0, my - (t + h))
            return (dx * dx + dy * dy) ** 0.5

        pairs = sorted(
            ((edge_dist(mb, box), fi, mi)
             for mi, (mb, _mk) in enumerate(marks)
             for fi, box in need
             if edge_dist(mb, box) <= 0.9),          # 0.9인치 밖이면 남의 것
            key=lambda x: x[0])
        took_f, took_m = set(), set()
        for _d, fi, mi in pairs:
            if fi in took_f or mi in took_m:
                continue
            fields[fi]["mark"] = marks[mi][1]
            took_f.add(fi)
            took_m.add(mi)

        for i, f in enumerate(fields):
            f["key"] = f"{'box' if f['kind'] == '상자' else 'arw'}:{f['mark'] or i}"

        out.append({"no": si, "title": title or f"{si}번 구조도", "fields": fields})
    return out


def thumb_path(no: int):
    """그 구조도의 미리보기 사진 경로(없으면 None)."""
    p = os.path.join(THUMBS, f"{no:02d}.png")
    return p if os.path.exists(p) else None


# ── 만든 구조도를 요약본에 끼워 넣기 ──────────────────
_A = "{http://schemas.openxmlformats.org/drawingml/2006/main}"
_P = "{http://schemas.openxmlformats.org/presentationml/2006/main}"
_R = "{http://schemas.openxmlformats.org/officeDocument/2006/relationships}"


def insert_diagram(dst_slide, pptx_path: str, logo_below_in: float = 0.95):
    """구조도 PPT(한 장)의 도형을 **그대로** 요약본 슬라이드에 옮겨 붙인다.

    그림으로 바꾸지 않는 이유: 웹 서버에는 파워포인트가 없어 그림으로 못 바꾼다.
    도형째 옮기면 변환이 필요 없고, 확대해도 안 깨지고, 나중에 글자도 고칠 수 있다.

    반환: 옮겨 붙인 도형들을 묶은 그룹(위치·크기는 부르는 쪽에서 정한다).
           넣을 게 없으면 None.
    """
    from lxml import etree
    src = Presentation(pptx_path).slides[0]
    sw = Presentation(pptx_path).slide_width

    picks = []
    for sh in src.shapes:                      # 맨 바깥 도형만(그룹은 통째로)
        try:
            l, t = sh.left or 0, sh.top or 0
            w, h = sh.width or 0, sh.height or 0
        except Exception:
            continue
        if w <= 0 or h <= 0:
            continue
        if t < Emu(int(logo_below_in * EMU_IN)):
            continue                            # 로고·구조도 제목은 뺀다
        if w > sw * 0.95:
            continue                            # 슬라이드 전체를 덮는 액자
        picks.append((sh, l, t, w, h))
    if not picks:
        return None

    x0 = min(p[1] for p in picks)
    y0 = min(p[2] for p in picks)
    x1 = max(p[1] + p[3] for p in picks)
    y1 = max(p[2] + p[4] for p in picks)
    cx, cy = max(1, x1 - x0), max(1, y1 - y0)

    grp = etree.SubElement(dst_slide.shapes._spTree, _P + "grpSp")
    nv = etree.SubElement(grp, _P + "nvGrpSpPr")
    c = etree.SubElement(nv, _P + "cNvPr")
    c.set("id", "9001")
    c.set("name", "금융구조도")
    etree.SubElement(nv, _P + "cNvGrpSpPr")
    etree.SubElement(nv, _P + "nvPr")
    spPr = etree.SubElement(grp, _P + "grpSpPr")
    xf = etree.SubElement(spPr, _A + "xfrm")
    for tag, a, b in (("off", x0, y0), ("ext", cx, cy),
                      ("chOff", x0, y0), ("chExt", cx, cy)):
        e = etree.SubElement(xf, _A + tag)
        if tag in ("off", "chOff"):
            e.set("x", str(int(a)))
            e.set("y", str(int(b)))
        else:
            e.set("cx", str(int(a)))
            e.set("cy", str(int(b)))

    for sh, *_ in picks:
        el = deepcopy(sh._element)
        # 그림이 가리키는 관계 번호는 슬라이드마다 다르다 → 다시 맺어준다.
        for node in el.iter():
            for attr in ("embed", "link", "id"):
                k = _R + attr
                rid = node.get(k)
                if not rid:
                    continue
                try:
                    rel = src.part.rels[rid]
                except KeyError:
                    continue
                try:
                    node.set(k, dst_slide.part.relate_to(
                        rel.target_ref if rel.is_external else rel._target,
                        rel.reltype, is_external=rel.is_external))
                except Exception:
                    pass
        grp.append(el)
    return dst_slide.shapes[-1]


def scale_tables_in_group(grp, ratio: float, min_pt: float = 4.5):
    """그룹을 줄일 때 따라오지 않는 것들을 직접 줄인다.

    그룹을 줄이면 도형의 '자리와 크기'는 같이 줄지만, 다음 둘은 안 줄어든다.
      · **표(graphicFrame)** — 아예 크기가 안 변한다(OOXML 특성). 상자들이 겹친다.
      · **글자 크기** — 도형은 작아지는데 글씨는 그대로라 밖으로 넘친다.
    그래서 열 너비·행 높이·글꼴 크기를 우리가 직접 곱해준다.
    """
    if ratio <= 0 or abs(ratio - 1.0) < 0.01:
        return

    # ① 글상자 글자 크기 (도형 자체는 그룹이 알아서 줄여준다)
    for sh, _p in _walk(grp.shapes):
        if sh.has_table or not sh.has_text_frame:
            continue
        for para in sh.text_frame.paragraphs:
            for run in para.runs:
                cur = run.font.size.pt if run.font.size else 11.0
                run.font.size = Pt(max(min_pt, round(cur * ratio, 1)))

    # ② 표 — 크기·글자 둘 다
    for sh, _p in _walk(grp.shapes):
        if not sh.has_table:
            continue
        t = sh.table
        try:
            for col in t.columns:
                col.width = int((col.width or 0) * ratio)
            for row in t.rows:
                row.height = int((row.height or 0) * ratio)
        except Exception:
            pass
        for row in t.rows:
            for cell in row.cells:
                for para in cell.text_frame.paragraphs:
                    for run in para.runs:
                        cur = run.font.size.pt if run.font.size else 9.0
                        run.font.size = Pt(max(min_pt, round(cur * ratio, 1)))
        try:                     # 칸 안쪽 여백도 같이 줄여야 글이 안 넘친다
            for row in t.rows:
                for cell in row.cells:
                    for a in ("margin_left", "margin_right",
                              "margin_top", "margin_bottom"):
                        v = getattr(cell, a, None)
                        if v:
                            setattr(cell, a, int(v * ratio))
        except Exception:
            pass


# ── 채운 구조도 만들기 ───────────────────────────────
def _set_text(tf, s):
    """서식(글꼴·크기·색)은 두고 글자만 바꾼다."""
    from engine_bits import replace_text_keep_runs
    replace_text_keep_runs(tf, s)


def _drop(sh):
    try:
        sh._element.getparent().remove(sh._element)
    except Exception:
        pass


def build_diagram(no: int, values: dict, out_path: str, removed=None,
                  layout: str = None) -> str:
    """고른 구조도(no)에 값을 채워 **한 장짜리 PPTX** 로 저장한다.

    values  : {field["key"]: "넣을 글자"} — 빈 문자열이면 글자만 지운다(상자는 남음).
              아예 빠진 key 는 틀의 글자를 그대로 둔다.
    removed : 아예 빼버릴 key 들. 글자만이 아니라 **상자·글상자째** 없앤다.
              (빈 글자만 지우면 빈 상자가 덩그러니 남아 이상해 보인다)
    번호 표식(①, ①-1 …)은 여기서 전부 지운다(회사 원본에는 없는 것).
    """
    removed = set(removed or ())
    prs = Presentation(layout or LAYOUT)
    keep = prs.slides[no - 1]

    # 고른 것만 남기고 나머지 슬라이드는 버린다
    ids = list(prs.slides._sldIdLst)
    for i, sid in enumerate(ids):
        if i != no - 1:
            prs.slides._sldIdLst.remove(sid)

    info = read_layouts(layout or LAYOUT)[no - 1]

    # ★반드시 '값 채우기' 를 먼저, '빼기' 를 나중에.
    #   도형을 먼저 지우면 뒤 도형들의 순번이 한 칸씩 밀려서, 값이 엉뚱한 칸에
    #   들어간다(시공사 자리에 신탁사 이름이 들어가던 문제).
    for f in info["fields"]:
        if f["key"] in removed or f["key"] not in values:
            continue
        val = (values.get(f["key"]) or "").strip()
        kind, path = f["loc"][0], f["loc"][1]
        sh = _by_path(keep, path)
        if sh is None:
            continue
        if kind == "cell":
            _, _, r, c = f["loc"]
            try:
                _set_text(sh.table.cell(r, c).text_frame, val)
            except Exception:
                pass
        else:
            _set_text(sh.text_frame, val)

    # 뺀 항목 없애기 — 뒤에서부터 지워야 남은 것들의 순번이 안 밀린다.
    for f in sorted((f for f in info["fields"] if f["key"] in removed),
                    key=lambda f: f["loc"][1], reverse=True):
        sh = _by_path(keep, f["loc"][1])
        if sh is None:
            continue
        if f["loc"][0] == "tb":
            _drop(sh)
            continue
        _, _, r, _c = f["loc"]
        try:
            t = sh.table
            if len(t.rows) <= 1:            # 한 줄짜리 상자면 상자째
                _drop(sh)
            else:                            # 여러 줄이면 그 줄만
                tr = t._tbl.tr_lst[r]
                tr.getparent().remove(tr)
        except Exception:
            _drop(sh)

    # 번호 표식 지우기 — 글상자째 없애거나, 글 앞머리에서 떼어낸다
    for sh, _p in list(_walk(keep.shapes)):
        if not sh.has_text_frame:
            continue
        txt = (sh.text_frame.text or "").strip()
        if not txt:
            continue
        if _MARK_ONLY.match(txt):
            _drop(sh)
            continue
        m = _MARK_HEAD.match(txt)
        if m and len(txt) > m.end():
            _set_text(sh.text_frame, txt[m.end():].strip())
    for sh, _p in list(_walk(keep.shapes)):
        if not sh.has_table:
            continue
        t = sh.table
        for r in range(len(t.rows)):
            for c in range(len(t.columns)):
                txt = (t.cell(r, c).text or "").strip()
                m = _MARK_HEAD.match(txt) if txt else None
                if m and len(txt) > m.end():
                    _set_text(t.cell(r, c).text_frame, txt[m.end():].strip())

    prs.save(out_path)
    return out_path
