# -*- coding: utf-8 -*-
"""요약본 파이프라인 — 원본 PDF → 4~5장 요약본.
  extract_summary(pdf_bytes) -> data(dict)
  build_summary(data, pdf_path, out_path, pages=1) -> out_path
"""
import io, os, json, sys
from copy import deepcopy

# ── 경로 ─────────────────────────────────────────────
#   생성기 코드는 'IM 자동화\생성기\' 에 있고,
#   변환 엔진(extractors.py, modules/)은 옆 저장소 'rainfield-im' 에 있다.
# ★요약본 생성기는 'IM 자동화' 폴더 안에서 **독립적으로** 돈다.
#   전엔 옆 저장소(rainfield-im)의 extractors·modules 를 끌어다 썼는데,
#   그 폴더가 움직이거나 정리될 때마다 깨졌다(경로·API키·캐시 3번 터짐).
#   → 필요한 것(extractors.py, claude_api.py, .env, 캐시)을 이 폴더에 두고 그것만 쓴다.
_HERE = os.path.dirname(os.path.abspath(__file__))        # ...\IM 자동화\_내부\생성기
_WORK = os.path.abspath(os.path.join(_HERE, ".."))        # ...\IM 자동화\_내부
if _HERE not in sys.path:
    sys.path.insert(0, _HERE)

from dotenv import load_dotenv
load_dotenv(os.path.join(_HERE, ".env"))                  # 로컬: 같은 폴더 .env
if not os.environ.get("ANTHROPIC_API_KEY"):
    # 배포(Streamlit Cloud): .env 가 없으므로 secrets 에서 받는다.
    try:
        import streamlit as _st
        _k = _st.secrets.get("ANTHROPIC_API_KEY", "")
        if _k:
            os.environ["ANTHROPIC_API_KEY"] = _k
    except Exception:
        pass

from pptx import Presentation
from pptx.util import Inches
from pptx.enum.text import PP_ALIGN, MSO_AUTO_SIZE
from extractors import extract_from_pdf
import claude_api as _claude_api
from claude_api import call_claude
from engine_bits import (replace_text_keep_runs as _replace_text_keep_runs,
                         extract_page_images, repeated_xrefs, mean_brightness)

from pathlib import Path as _Path
_claude_api.CACHE_DIR = _Path(_HERE) / ".claude_cache"    # 캐시도 이 폴더에 고정
_claude_api.CACHE_DIR.mkdir(parents=True, exist_ok=True)

# 레이아웃 틀 — 로컬(IM 자동화\_내부\레이아웃\)과 배포(생성기\layout\) 둘 다 지원
_LAYOUT_CANDIDATES = [
    os.path.join(_HERE, "layout", "요약본 레이아웃.pptx"),   # 저장소 배포용
    os.path.join(_WORK, "레이아웃", "요약본 레이아웃.pptx"),  # 로컬
]


def _layout_path():
    for p in _LAYOUT_CANDIDATES:
        if os.path.exists(p):
            return p
    raise FileNotFoundError(
        "레이아웃 틀(요약본 레이아웃.pptx)을 찾을 수 없습니다.\n" +
        "\n".join(f"  - {p}" for p in _LAYOUT_CANDIDATES))


SYS = """당신은 부동산 금융 IM(투자설명서)을 읽고 '요약본 제안서'용 핵심 데이터를 뽑는 전문가다.
아래 원문에서 정보를 찾아 **JSON만** 출력한다. 규칙:
- 값은 원문 그대로(verbatim). 숫자·단위·괄호([30]억원, [43.98]% 등 대괄호 표기)도 원문대로 유지.
- 원문에 없는 값은 지어내지 말고 null. 확인 불가도 null.
- '본건'에 해당하는 트랜치(Tr)는 is_bongeon=true.

출력 JSON 스키마:
{
 "deal_name": "사업명(표지용)",
 "date_ko": "발행일 또는 제안 시점(예: 2026년 8월)",
 "highlights": [
   {"title":"12자내 짧은 제목", "subtitle":"20자내 한 줄 부제", "bullets":["근거 1(상세)","근거 2(상세)"]}
 ],
 "사모사채개요": {"사모사채명":..,"사채유형":..,"발행인":..,"기초자산":..,"발행금액":..,"발행일":..,"만기일":..,"금융조건":..,"이자지급주기":..},
 "담보대출조건": {"금융구조":..,"차주":..,"tranches":[{"구분":"Tr.A","대출금액":..,"LTV":..,"금리":..,"is_bongeon":false}],"LTV기준_각주":..,"주요채권보전":["..",".."],"주요인출선행조건":["..",".."],"상환재원":..,"대출기간":..},
 "사업일정": {"완료":[{"시기":"'23.02","내용":".."}],"예정":[{"시기":"'26.4Q","내용":".."}]},
 "사업개요": {"사업명":..,"대지위치":..,"대지면적":..,"건축면적":..,"연면적":..,"건축규모":..},
 "투입에쿼티": [{"주체":..,"금액":..,"비고":..}],
 "법인개요": {"대상":..,"회사명":..,"설립일":..,"사업자번호":..,"대표자":..,"자본금":..,"주요사업":..,"주요주주":..},
 "재무제표": {"단위":"백만원","연도":["2025.12"],"행":[{"항목":"자산총계","값":[".."]}]},
 "중요항목": [{"제목":..,"내용":["..",".."]}],
 "이미지_있음": {"금융구조도": false, "조감도": true, "위치도": false}
}
사업일정은 핵심 마일스톤만(완료+예정 합쳐 6~8개 이내). JSON 외 다른 말 금지."""


def extract_summary(pdf_bytes: bytes) -> dict:
    data = extract_from_pdf(pdf_bytes)
    full = "\n".join(data.get("pages_text", []))
    res = call_claude(SYS, f"[IM 원문]\n{full}", slide_num=801,
                      pdf_context=full, prompt_version="summary_extract_v2")
    if not res.get("ok"):
        raise RuntimeError(res.get("error") or "추출 실패")
    return res["data"]


# ── 빌드 헬퍼 ────────────────────────────────────────
def _walk(shapes):
    for sh in shapes:
        yield sh
        if sh.shape_type == 6:
            yield from _walk(sh.shapes)


def _set(table, r, c, text):
    if r < len(table.rows) and c < len(table.columns):
        _replace_text_keep_runs(table.cell(r, c).text_frame, str(text if text is not None else ""))


def _lines(table, r, c, lines):
    _replace_text_keep_runs(table.cell(r, c).text_frame, "\n".join(str(x) for x in lines if x))


def _add_rows(table, n):
    for _ in range(n):
        table._tbl.append(deepcopy(table._tbl.tr_lst[-1]))


def _del_rows(table, n):
    for _ in range(n):
        tr = table._tbl.tr_lst[-1]; tr.getparent().remove(tr)


def _replace_contains(slide, needle, new_text):
    for sh in _walk(slide.shapes):
        if sh.has_text_frame and needle in sh.text_frame.text:
            _replace_text_keep_runs(sh.text_frame, new_text)
            return True
    return False


def _find_label(slide, needle):
    """슬라이드에서 특정 문구가 들어간 '라벨 글상자'를 찾는다(표 제목용)."""
    for sh in slide.shapes:
        if sh.has_text_frame and needle in sh.text_frame.text:
            return sh
    return None


def _dup_slide_after(prs, src, after_idx):
    """슬라이드를 통째로 복제해 after_idx 바로 뒤에 넣는다(핵심요약 2페이지용).

    python-pptx 에 복제 기능이 없어 직접 한다:
      ①같은 레이아웃으로 새 슬라이드 → 자동 생성된 자리표시자 제거
      ②원본 도형 XML 을 deepcopy 해서 그대로 붙임
      ③그림·차트 같은 외부 파트 관계(rel)를 함께 복사 — 안 하면 이미지가 깨진다
      ④맨 뒤에 생긴 슬라이드를 원하는 위치로 이동
    """
    import copy as _copy
    dst = prs.slides.add_slide(src.slide_layout)
    for sh in list(dst.shapes):
        sh._element.getparent().remove(sh._element)
    for sh in src.shapes:
        dst.shapes._spTree.append(_copy.deepcopy(sh._element))
    for rid, rel in src.part.rels.items():
        try:
            if rel.is_external:
                dst.part.rels.add_relationship(rel.reltype, rel._target, rid, True)
            else:
                dst.part.rels.add_relationship(rel.reltype, rel._target, rid)
        except Exception:
            pass
    # 맨 뒤 → after_idx 다음 자리로 이동
    sldIdLst = prs.slides._sldIdLst
    ids = list(sldIdLst)
    sldIdLst.remove(ids[-1])
    sldIdLst.insert(after_idx + 1, ids[-1])
    return dst


def _table_h(shape):
    """표의 실제 높이(행 높이 합)."""
    try:
        return sum(r.height for r in shape.table.rows)
    except Exception:
        return shape.height or 0


def _stack(blocks, x_in, w_in, top_in, bottom_in, gap_in=0.16, label_gap_in=0.05):
    """한 단(column)의 블록들을 위에서부터 차곡차곡 다시 배치한다.

    ★레이아웃 틀은 위치가 '고정'이라, 표 행 수가 원본에 따라 늘면 아래 블록과 겹친다
      (연대보증사 재무제표가 법인개요를 덮던 문제). 내용을 채운 뒤 실제 높이로 재배치한다.
    blocks: [(label_shape|None, table_or_pic_shape|None), ...]  위→아래 순서
    반환: 마지막 y(인치). bottom 을 넘으면 넘친 것이므로 호출측에서 판단.
    """
    X, W = Inches(x_in), Inches(w_in)
    y = Inches(top_in)
    for label, body in blocks:
        if label is not None:
            label.left, label.top = X, y
            y = y + (label.height or Inches(0.25)) + Inches(label_gap_in)
        if body is not None:
            body.left, body.top, body.width = X, y, W
            h = _table_h(body) if body.has_table else (body.height or 0)
            body.height = h
            y = y + h + Inches(gap_in)
    return y / 914400.0


def build_highlight_preview(data: dict, out_path: str) -> str:
    """하이라이트 슬라이드 1장만 만든 PPTX(미리보기용)."""
    return build_summary(data, None, out_path, pages=1, _highlight_only=True)


# ── 하이라이트 제목 블록([체크아이콘]+[제목]) 가운데 정렬 ───────────────
#   문제: 레이아웃 틀의 제목 그룹 3개가 전부 L4.66" W1.52" 로 **고정**이고
#   그 안 글상자는 폭 1.04" 짜리다. 제목이 그보다 길면 글이 상자를 뚫고
#   오른쪽으로만 자라서, **글자 수가 다르면 카드마다 제목 위치가 어긋났다.**
#   (예: '초역세권 입지' 6자 vs '낮은 인허가 리스크' 9자)
#   해결: 제목 글자폭을 재서 그룹 폭을 실제 내용에 맞추고, 그룹 전체를
#   슬라이드 정중앙에 놓는다. 기존 IM 변환기의 하이라이트와 같은 방식.
_TITLE_PT = 18.0    # 틀의 제목 rPr 에 sz 가 없다 → 본문 기본값 18pt
_ICON_GAP = 0.08    # 체크아이콘과 제목 사이 간격(인치)
_A = "{http://schemas.openxmlformats.org/drawingml/2006/main}"


def _text_width_in(s: str, pt: float) -> float:
    """글자폭 추정(인치). 한글·한자는 전각(1.0em), 숫자·영문·기호는 반각(0.55em)."""
    u = 0.0
    for ch in s:
        if ('가' <= ch <= '힣' or '㄰' <= ch <= '㆏'
                or '一' <= ch <= '鿿' or ch in '·／…‘’“”'):
            u += 1.0
        elif ch == ' ':
            u += 0.30
        else:
            u += 0.55
    return u * pt / 72.0


def _center_title_group(grp, txbox, title: str, slide_w: int) -> bool:
    """[체크아이콘 + 제목] 그룹을 글자폭에 맞춰 다시 재고 슬라이드 가운데로 옮긴다."""
    xfrm = grp._element.find(f"{_A}xfrm")
    if xfrm is None:                       # grpSpPr 안에 있는 경우
        xfrm = grp._element.find(f".//{_A}xfrm")
    if xfrm is None:
        return False
    off, ext = xfrm.find(f"{_A}off"), xfrm.find(f"{_A}ext")
    ch_off, ch_ext = xfrm.find(f"{_A}chOff"), xfrm.find(f"{_A}chExt")
    if None in (off, ext, ch_off, ch_ext):
        return False

    pics = [c for c in grp.shapes if c.shape_type == 13]
    if not pics:
        return False
    icon = pics[0]
    icon_w = icon.width or Inches(0.47)

    base_x = int(ch_off.get("x"))          # 자식 좌표계의 왼쪽 끝
    text_w = int(Inches(_text_width_in(title or "", _TITLE_PT) + 0.06))
    text_w = max(text_w, Inches(0.4))

    # 자식 좌표계 재배치: [아이콘][간격][제목]
    icon.left = base_x
    txbox.left = base_x + icon_w + Inches(_ICON_GAP)
    txbox.width = text_w

    total = icon_w + Inches(_ICON_GAP) + text_w
    ext.set("cx", str(int(total)))          # 확대·축소 없이 1:1 이 되도록
    ch_ext.set("cx", str(int(total)))
    off.set("x", str(int((slide_w - total) / 2)))   # 슬라이드 정중앙
    return True


def build_summary(data: dict, pdf_path: str, out_path: str, pages: int = 1,
                  _highlight_only: bool = False) -> str:
    prs = Presentation(_layout_path())

    # 슬라이드 1 : 표지
    s1 = prs.slides[0]
    name = (data.get('deal_name') or '').strip()
    _replace_contains(s1, "사업명(요약본)", f"{name} 사모사채 제안서(요약본)")
    _replace_contains(s1, "몇월", data.get('date_ko', ''))

    # 슬라이드 2 : 하이라이트
    s2 = prs.slides[1]
    hl = data.get('highlights') or []
    title_boxes = []
    for sh in s2.shapes:
        if sh.shape_type == 6:
            for ch in sh.shapes:
                if ch.has_text_frame and ch.text_frame.text.strip() == "제목":
                    title_boxes.append((sh.top, sh, ch))   # (정렬키, 그룹, 글상자)
    title_boxes = [(g, c) for _, g, c in sorted(title_boxes, key=lambda x: x[0])]
    sub_boxes = sorted([sh for sh in s2.shapes if sh.has_text_frame and sh.text_frame.text.strip() == "부제목"],
                       key=lambda sh: sh.top)
    bull_boxes = sorted([sh for sh in s2.shapes if sh.has_text_frame and len(sh.text_frame.text.strip()) > 20],
                        key=lambda sh: sh.top)
    # 슬라이드 폭(요약본 레이아웃은 10.83in) 기준으로 꺾쇠 안쪽에 맞춘다
    _SW = prs.slide_width
    # 꺾쇠는 1.45~9.39in 이지만 양끝의 화살촉이 안쪽으로 파고든다.
    # 7.40in 로 두면 긴 줄이 오른쪽 화살촉에 닿았다 → 기존 IM 변환기와 같은 여백으로.
    _BODY_L, _BODY_W = Inches(1.80), Inches(7.00)

    for i in range(min(3, len(hl))):
        if i < len(title_boxes):
            grp, tb = title_boxes[i]
            _title = (hl[i].get('title') or '').strip()
            _replace_text_keep_runs(tb.text_frame, _title)
            # 제목은 아이콘 오른쪽에서 시작(왼쪽정렬)하고, 아이콘까지 묶은 블록을
            # _center_title_group 이 슬라이드 정중앙으로 옮긴다.
            tb.text_frame.word_wrap = False
            for p in tb.text_frame.paragraphs:
                p.alignment = PP_ALIGN.LEFT
            _center_title_group(grp, tb, _title, _SW)
        if i < len(sub_boxes):
            sb = sub_boxes[i]
            _replace_text_keep_runs(sb.text_frame, hl[i].get('subtitle', ''))
            # 부제목은 가운데정렬이라 폭만 넉넉히 주면 알아서 가운데에 선다.
            # (틀의 2.58" 를 그대로 두면 긴 부제목이 두 줄로 접힌다.)
            sb.left, sb.width = _BODY_L, _BODY_W
            sb.text_frame.word_wrap = False
            for p in sb.text_frame.paragraphs:
                p.alignment = PP_ALIGN.CENTER
        if i < len(bull_boxes):
            bb = bull_boxes[i]
            # 틀이 이미 '→' 를 글머리기호로 찍는다. 본문에도 '→' 가 들어오면
            # '→ →' 로 두 번 찍히므로 앞머리의 기호를 떼어낸다.
            _lines = []
            for _b in (hl[i].get('bullets') or []):
                _b = str(_b).strip()
                while _b[:1] in ("→", "·", "-", "•", "▶", "ㆍ"):
                    _b = _b[1:].strip()
                if _b:
                    _lines.append(_b)
            _replace_text_keep_runs(bb.text_frame, "\n".join(_lines))
            # ★근거: word_wrap 이 꺼져 있어 한 줄이 꺾쇠 밖으로 뚫고 나갔다 → 켜서 안에서 접히게.
            bb.left, bb.width = _BODY_L, _BODY_W
            bb.text_frame.word_wrap = True
            try:
                bb.text_frame.auto_size = MSO_AUTO_SIZE.NONE
            except Exception:
                pass

    if _highlight_only:
        # 미리보기: 하이라이트 슬라이드만 남기고 나머지 삭제
        _ids = list(prs.slides._sldIdLst)
        for _i, _sid in enumerate(_ids):
            if _i != 1:
                prs.slides._sldIdLst.remove(_sid)
        prs.save(out_path)
        return out_path

    # 슬라이드 3 : 핵심 요약
    s3 = prs.slides[2]
    # ★shape.table 은 접근할 때마다 새 객체를 만들어 'is' 비교가 안 된다
    #   → 나중에 재배치할 수 있도록 '도형' 참조도 같은 순서로 잡아둔다.
    tshapes = [sh for sh in s3.shapes if sh.has_table]
    tables = [sh.table for sh in tshapes]
    T_LOAN, T_TR, T_FIN, T_CORP, T_BOND, T_SCHED = tables
    S_LOAN, S_TR, S_FIN, S_CORP, S_BOND, S_SCHED = tshapes
    loan = data.get('담보대출조건') or {}
    bong = next((t for t in (loan.get('tranches') or []) if t.get('is_bongeon')), {})

    # 사모사채 개요(고정+파생)
    기간 = (loan.get('대출기간') or '').strip()
    for r, v in {3: f"{name} 담보대출채권", 4: bong.get('대출금액'),
                 5: f"{data.get('date_ko','')} (예정)",
                 6: f"발행일로부터 {기간 or '약 3개월'} (대출 만기 연동)",
                 7: f"채권 금리 : 연 {bong.get('금리','')} (고정, 세전)"}.items():
        if v and r < len(T_BOND.rows):
            _set(T_BOND, r, 1, v)

    # 담보대출
    _set(T_LOAN, 0, 1, loan.get('금융구조'))
    _set(T_LOAN, 1, 1, loan.get('차주'))
    if len(T_LOAN.rows) > 2: _set(T_LOAN, 2, 1, loan.get('LTV기준_각주'))
    if len(T_LOAN.rows) > 3: _lines(T_LOAN, 3, 1, loan.get('주요채권보전') or [])
    if len(T_LOAN.rows) > 4: _set(T_LOAN, 4, 1, loan.get('대출기간'))

    # Tr
    trs = loan.get('tranches') or []
    ncol = len(T_TR.columns)
    _set(T_TR, 0, 0, "구분")
    for c in range(1, ncol):
        if c-1 < len(trs):
            _set(T_TR, 0, c, trs[c-1].get('구분', '') + (" (본건)" if trs[c-1].get('is_bongeon') else ""))
    if len(T_TR.rows) >= 2:
        _set(T_TR, 1, 0, "대출금액")
        for c in range(1, ncol):
            if c-1 < len(trs): _set(T_TR, 1, c, trs[c-1].get('대출금액'))
    if len(T_TR.rows) >= 3:
        _set(T_TR, 2, 0, "LTV / 금리")
        for c in range(1, ncol):
            if c-1 < len(trs):
                v = trs[c-1]; _set(T_TR, 2, c, f"{v.get('LTV') or '-'} / {v.get('금리') or '-'}")

    # 재무
    fin = data.get('재무제표') or {}
    if fin:
        yrs = fin.get('연도') or []; rows = fin.get('행') or []
        _set(T_FIN, 0, 0, f"구분\n(단위 : {fin.get('단위','백만원')})")
        for c in range(1, len(T_FIN.columns)):
            _set(T_FIN, 0, c, yrs[c-1] if c-1 < len(yrs) else "")
        need = len(rows) + 1; cur = len(T_FIN.rows)
        if need > cur: _add_rows(T_FIN, need-cur)
        elif need < cur: _del_rows(T_FIN, cur-need)
        for i, row in enumerate(rows, start=1):
            _set(T_FIN, i, 0, row.get('항목'))
            for c in range(1, len(T_FIN.columns)):
                vals = row.get('값') or []
                _set(T_FIN, i, c, vals[c-1] if c-1 < len(vals) else "")

    # 법인개요
    corp = data.get('법인개요') or {}
    if corp:
        pairs = [("회사명", corp.get('회사명'), "사업자번호", corp.get('사업자번호')),
                 ("설립일", corp.get('설립일'), "대표자", corp.get('대표자')),
                 ("자본금", corp.get('자본금'), "주요사업", corp.get('주요사업')),
                 ("주요주주", corp.get('주요주주'), "", "")]
        need = len(pairs); cur = len(T_CORP.rows)
        if need < cur: _del_rows(T_CORP, cur-need)
        elif need > cur: _add_rows(T_CORP, need-cur)
        for i, (a, b, c, d) in enumerate(pairs):
            _set(T_CORP, i, 0, a); _set(T_CORP, i, 1, b)
            if len(T_CORP.columns) > 2:
                _set(T_CORP, i, 2, c); _set(T_CORP, i, 3, d)

    # 사업일정(핵심만)
    sched = data.get('사업일정') or {}
    done = sched.get('완료') or []; plan = sched.get('예정') or []
    allrows = [("완료", d) for d in done] + [("예정", p) for p in plan]
    need = len(allrows); cur = len(T_SCHED.rows)
    if need > cur: _add_rows(T_SCHED, need-cur)
    elif need < cur: _del_rows(T_SCHED, cur-need)
    for i, (grp, it) in enumerate(allrows):
        _set(T_SCHED, i, 0, grp if (i == 0 or i == len(done)) else "")
        _set(T_SCHED, i, 1, it.get('시기')); _set(T_SCHED, i, 2, it.get('내용'))

    # ── 라벨을 이 딜의 회사명으로 교체 ─────────────────
    #   틀에는 참고자료(신영증권 건)의 '(주)신영' 이 박혀 있어, 그대로 두면
    #   남의 회사 이름이 우리 제안서에 나간다. 반드시 실제 회사명으로 바꾼다.
    _corp_nm = (corp.get('회사명') or data.get('연대보증사') or '').strip()
    lbl_corp = _find_label(s3, "법인개요")
    lbl_fin = _find_label(s3, "재무제표")
    if lbl_corp is not None:
        _replace_text_keep_runs(lbl_corp.text_frame,
                                f"{_corp_nm} 법인개요" if _corp_nm else "법인개요")
    if lbl_fin is not None:
        _replace_text_keep_runs(lbl_fin.text_frame,
                                f"{_corp_nm} 주요 재무현황" if _corp_nm else "주요 재무현황")
    lbl_loan = _find_label(s3, "담보대출")
    lbl_sched = _find_label(s3, "사업일정")
    lbl_diag = _find_label(s3, "금융구조도")

    # ── 참고자료 이미지 제거 + 이 딜의 조감도 삽입 ──────
    diag_pic = None
    for sh in list(s3.shapes):
        if sh.shape_type == 13 and sh.width and sh.width > int(Inches(1.8)):
            if diag_pic is None and (sh.top or 0) > int(Inches(4.5)):
                diag_pic = sh          # 좌하단 = 금융구조도 자리(살려서 재배치)
                continue
            sh._element.getparent().remove(sh._element)
        elif sh.shape_type == 6:
            # ★참고 레이아웃의 '샘플 조감도'는 GROUP 안에 그림 2장으로 들어있다.
            #   PICTURE 조건만 보면 GROUP 안은 못 잡아 그대로 남고, 슬라이드를 복제하면
            #   그 그림들이 '그림을 표시할 수 없습니다'로 깨져 나온다. → 그림 든 GROUP 은 제거.
            try:
                has_pic = any(c.shape_type == 13 for c in sh.shapes)
            except Exception:
                has_pic = False
            if has_pic or (sh.top and sh.top > int(Inches(5.5))):
                sh._element.getparent().remove(sh._element)
    if diag_pic is not None:           # 아직 작도 전 → 자리만 두면 지저분하므로 제거
        diag_pic._element.getparent().remove(diag_pic._element)
        if lbl_diag is not None:
            lbl_diag._element.getparent().remove(lbl_diag._element)
            lbl_diag = None

    top_pic = None
    try:
        import fitz
        doc = fitz.open(pdf_path)
        # ★로고·워터마크(여러 페이지 반복) 제외 — 안 그러면 발행사 로고가 조감도 자리에 박힌다
        skip = repeated_xrefs(doc)
        imgs = []
        for pi in range(doc.page_count):
            imgs += extract_page_images(doc, pi, skip_xrefs=skip)
        imgs = [im for im in imgs if im.get('data')]
        # ★조감도 고르기: ①사진(평균밝기<200 — 공문·도면은 흰 바탕이라 밝다)
        #   ②가로로 넓은 것 ③큰 것.  전엔 '가장 큰 것'만 봐서 삼성증권 공문이 뽑혔다.
        for im in imgs:
            ar = im['width'] / max(1, im['height'])
            im['_score'] = (1 if mean_brightness(im['data']) < 200 else 0,
                            1 if 1.2 <= ar <= 2.6 else 0,
                            im['width'] * im['height'])
        imgs.sort(key=lambda im: im['_score'], reverse=True)
        if imgs:
            top_pic = s3.shapes.add_picture(
                io.BytesIO(imgs[0]['data']), Inches(5.26), Inches(0.39),
                width=Inches(5.32), height=Inches(1.80))
        doc.close()
    except Exception as _ie:
        print(f"[요약본] 조감도 삽입 실패: {_ie}")

    # ── ★사용자가 3단계에서 정한 '칸 배치'대로 쌓는다 ──────
    #   data["_slots"] = {"left": [항목,...], "right": [항목,...]}  (위→아래 순서)
    #   없으면 기본 배치를 쓴다. 고른 것만 남기고 나머지 도형은 지운다.
    #   틀은 위치가 고정이라 표 행 수가 딜마다 달라지면 아래를 덮으므로,
    #   내용을 채운 뒤 실제 높이로 다시 쌓는다.
    DEFAULT1 = {"left":  ["사모사채개요", "담보대출조건"],
                "right": ["조감도", "사업일정", "법인개요", "재무제표"]}
    DEFAULT2 = {"left":  ["법인개요", "재무제표"], "right": []}
    slots = dict(data.get("_slots") or DEFAULT1)
    slots2 = dict(data.get("_slots2") or (DEFAULT2 if pages >= 2 else {}))

    def _blocks_of(slide):
        """그 슬라이드의 '항목명 → (라벨, 본문)' 지도.
           복제본도 표 순서·라벨 문구가 같으므로 똑같이 찾을 수 있다."""
        ts = [sh for sh in slide.shapes if sh.has_table]
        L, TR, F, C, B, SC = (ts + [None] * 6)[:6]
        pic = next((sh for sh in slide.shapes
                    if sh.shape_type == 13 and (sh.width or 0) > Inches(4)
                    and (sh.top or 0) < Inches(2.6)), None)
        return {
            "사모사채개요": (None, B),
            "담보대출조건": (_find_label(slide, "담보대출"), L),
            "대출조건표":   (None, TR),
            "사업일정":     (_find_label(slide, "사업일정"), SC),
            "법인개요":     (_find_label(slide, "법인개요"), C),
            "재무제표":     (_find_label(slide, "재무현황"), F),
            "조감도":       (None, pic),
        }

    def _arrange(slide, sl, tag):
        """고른 블록만 남기고 위에서부터 다시 쌓는다."""
        # ★칸은 사용자가 고른 그대로 쓴다. 임의로 묶거나 끼워 넣지 않는다
        #   (대출조건표도 별개 항목으로 원하는 자리에 배치 — 자동으로 붙이지 않음).
        BL = _blocks_of(slide)
        used = {k for side in ("left", "right") for k in sl.get(side, [])}
        for key, (lb, bd) in BL.items():
            if key in used:
                continue
            for sh in (lb, bd):
                if sh is not None:
                    try:
                        sh._element.getparent().remove(sh._element)
                    except Exception:
                        pass
        _BOT = 7.28
        for side, x, w, nm in (("left", 0.25, 4.90, "좌단"), ("right", 5.26, 5.32, "우단")):
            seq = [BL[k] for k in sl.get(side, []) if k in BL]
            seq = [(lb, bd) for lb, bd in seq if bd is not None or lb is not None]
            end = _stack(seq, x_in=x, w_in=w, top_in=0.39, bottom_in=_BOT)
            if end > _BOT:
                print(f"[요약본] 경고: {tag} {nm} 내용이 {end - _BOT:.2f}in 넘침 "
                      f"— 칸을 줄이거나 페이지를 나눠야 함")

    if pages >= 2:
        # ★내용을 다 채운 뒤 통째로 복제 → 각 장에서 필요없는 블록만 지운다.
        #   (빈 슬라이드를 새로 만들어 채우면 서식·표 스타일을 다시 맞춰야 해서 위험)
        s3b = _dup_slide_after(prs, s3, 2)
        # ★복제본의 그림은 관계(rel)가 원본 파트를 가리켜 '그림을 표시할 수 없습니다'로 깨진다.
        #   둘째 장에서 조감도를 안 쓰면 그냥 지운다(쓰는 경우만 남기고 아래에서 재삽입).
        _keep_pic2 = "조감도" in (list(slots2.get("left") or []) +
                                  list(slots2.get("right") or []))
        for sh in list(s3b.shapes):
            if sh.shape_type == 13 and (sh.width or 0) > Inches(1.8) and not _keep_pic2:
                sh._element.getparent().remove(sh._element)
        _arrange(s3, slots, "핵심요약 1/2")
        _arrange(s3b, slots2, "핵심요약 2/2")
    else:
        _arrange(s3, slots, "핵심요약")

    prs.save(out_path)
    return out_path
